from io import BytesIO
from typing import Iterable

from pptx import Presentation

from models.generation import PresentationSlide


def build_presentation(title: str, slides: Iterable[PresentationSlide]) -> BytesIO:
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title

    for slide_data in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_data.title
        slide.placeholders[1].text = slide_data.content

    file = BytesIO()
    presentation.save(file)
    file.seek(0)
    return file
